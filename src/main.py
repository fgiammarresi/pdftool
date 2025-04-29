import argparse
import os
import logging
from deep_translator import GoogleTranslator
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from tqdm import tqdm
import fitz  # PyMuPDF
from transformers import pipeline

# Imposta la configurazione del logger
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Funzione per la trascrizione
def trascrizione(file_path, output_format="txt"):
    logging.info(f"[TRASCRIZIONE] Elaboro il file: {file_path}")
    doc = fitz.open(file_path)
    pages_text = [page.get_text() for page in doc]
    doc.close()

    output_format = output_format.lower()

    if output_format == "txt":
        output_path = "trascrizione.txt"
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n\n".join(pages_text))

    elif output_format == "docx":
        output_path = "trascrizione.docx"
        document = Document()
        for text in pages_text:
            document.add_paragraph(text)
            document.add_page_break()
        document.save(output_path)

    elif output_format == "xlsx":
        output_path = "trascrizione.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Testo"
        for i, text in enumerate(pages_text, start=1):
            ws.cell(row=i, column=1, value=text)
        wb.save(output_path)

    elif output_format == "pptx":
        output_path = "trascrizione.pptx"
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[1]  # titolo + contenuto
        for i, text in enumerate(pages_text, start=1):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.title.text = f"Pagina {i}"
            slide.placeholders[1].text = text
        prs.save(output_path)

    else:
        logging.error(f"❌ Formato di output non supportato: {output_format}")
        return

    logging.info(f"✅ Trascrizione completata. File salvato come: {output_path}")

# Funzione per la traduzione
def traduzione(file_path, language="en"):
    logging.info(f"[TRADUZIONE] Elaboro il file: {file_path} e traducendo in {language}")
    doc = fitz.open(file_path)
    pages_text = [page.get_text() for page in doc]
    doc.close()

    translator = GoogleTranslator(target=language)
    translated_text = [translator.translate(page_text) for page_text in pages_text]

    output_path = f"traduzione_{language}.txt"
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(translated_text))

    logging.info(f"✅ Traduzione completata. File salvato come: {output_path}")

# Funzione per la sintesi
def sintesi(file_path):
    logging.info(f"[SINTESI] Elaboro il file: {file_path}")
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    doc = fitz.open(file_path)
    pages_text = [page.get_text() for page in doc]
    doc.close()

    text_to_summarize = " ".join(pages_text)
    summary = summarizer(text_to_summarize, max_length=150, min_length=50, do_sample=False)

    output_path = "sintesi.txt"
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(summary[0]['summary_text'])

    logging.info(f"✅ Sintesi completata. File salvato come: {output_path}")

# Funzione di analisi degli argomenti della riga di comando
def parse_args():
    parser = argparse.ArgumentParser(description="PDF Tool - Trascrizione, Traduzione, Sintesi")
    parser.add_argument('--action', type=str, required=True, help="Azione da eseguire: transcribe, translate, summarize")
    parser.add_argument('--file', type=str, required=True, help="Percorso del file PDF")
    parser.add_argument('--output', type=str, default="txt", help="Formato di output: txt, docx, xlsx, pptx")
    parser.add_argument('--lang', type=str, default="en", help="Lingua di destinazione per la traduzione")
    return parser.parse_args()

def main():
    # Analizza gli argomenti
    args = parse_args()

    # Gestione dei file con spazi (aggiungi virgolette per evitare problemi di separazione degli argomenti)
    file_path = args.file
    if " " in file_path:
        file_path = f'"{file_path}"'

    # Verifica se il file esiste
    if not os.path.exists(file_path.strip('"')):
        logging.error(f"Errore: il file {file_path} non esiste.")
        return

    # Esegui l'azione richiesta
    if args.action == "transcribe":
        trascrizione(args.file, args.output)
    elif args.action == "translate":
        traduzione(file_path, args.lang)
    elif args.action == "summarize":
        sintesi(file_path)
    else:
        logging.error(f"Azione non valida: {args.action}")

if __name__ == "__main__":
    main()
